from flask import Flask, request, send_file
from pdf2docx import Converter
import os
import uuid
import tempfile
import pdfplumber
import openpyxl
from openpyxl.styles import Font, Border, Side, PatternFill
from openpyxl.utils import get_column_letter
import fitz
import gc
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import io


app = Flask(__name__)


@app.route("/health", methods=["GET"])
def health():
    return {"status": "up"}


@app.route("/convert", methods=["POST"])
def convert():
    file = request.files["file"]
    temp_dir = tempfile.gettempdir()
    input_path = os.path.join(temp_dir, f"{uuid.uuid4()}.pdf")
    output_path = os.path.join(temp_dir, f"{uuid.uuid4()}.docx")

    file.save(input_path)

    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()

    os.remove(input_path)

    return send_file(output_path, as_attachment=True, download_name="converted.docx")


@app.route("/convert-to-excel", methods=["POST"])
def convert_to_excel():
    file = request.files["file"]
    temp_dir = tempfile.gettempdir()
    input_path = os.path.join(temp_dir, f"{uuid.uuid4()}.pdf")
    output_path = os.path.join(temp_dir, f"{uuid.uuid4()}.xlsx")

    file.save(input_path)

    workbook = openpyxl.Workbook()
    workbook.remove(workbook.active)

    def try_number(value):
        if value is None:
            return None
        cleaned = str(value).replace(",", "").strip()
        try:
            if "." in cleaned:
                return float(cleaned)
            return int(cleaned)
        except ValueError:
            return value

    thin = Side(style="thin", color="999999")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    header_fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")

    page_count = 0

    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            table_objects = page.find_tables()
            table_rows_list = page.extract_tables()

            if not table_objects or not table_rows_list:
                continue

            page_count += 1
            sheet = workbook.create_sheet(f"Page {page_count}")

            words = page.extract_words()

            table_line_texts = set()
            blocks = []
            table_bboxes = []

            paired = list(zip(table_objects, table_rows_list))

            for t_obj, rows in paired:
                for row in rows:
                    line_text = " ".join(str(c) for c in row if c).strip()
                    if line_text:
                        table_line_texts.add(line_text)
                blocks.append({
                    "type": "table",
                    "top": t_obj.bbox[1],
                    "rows": rows,
                })
                table_bboxes.append((t_obj.bbox[1], t_obj.bbox[3]))

            lines_by_top = {}
            for w in words:
                key = round(w["top"] / 3) * 3
                lines_by_top.setdefault(key, []).append(w["text"])

            for top, word_list in lines_by_top.items():
                line_text = " ".join(word_list).strip()
                if not line_text:
                    continue
                if line_text in table_line_texts:
                    continue

                inside_table = False
                for top_bound, bottom_bound in table_bboxes:
                    if top_bound - 2 <= top <= bottom_bound + 2:
                        inside_table = True
                        break

                if inside_table:
                    continue

                blocks.append({
                    "type": "text",
                    "top": top,
                    "text": line_text,
                })

            blocks.sort(key=lambda b: b["top"])

            current_row = 1
            width_tracker = {}
            max_cols_seen = 1
            heading_rows = []

            for block in blocks:
                if block["type"] == "text":
                    sheet.cell(row=current_row, column=1, value=block["text"])
                    sheet.cell(row=current_row, column=1).font = Font(bold=True, size=11)
                    heading_rows.append(current_row)
                    current_row += 1

                else:
                    rows = block["rows"]
                    header_row = current_row
                    num_cols = len(rows[0]) if rows else 1
                    max_cols_seen = max(max_cols_seen, num_cols)

                    for row in rows:
                        converted_row = [try_number(cell) for cell in row]
                        for col_index, value in enumerate(converted_row, start=1):
                            cell = sheet.cell(row=current_row, column=col_index, value=value)
                            cell.border = border

                            text_len = len(str(value)) if value is not None else 0
                            width_tracker[col_index] = max(width_tracker.get(col_index, 0), text_len)

                        current_row += 1

                    for col_index in range(1, num_cols + 1):
                        header_cell = sheet.cell(row=header_row, column=col_index)
                        header_cell.font = Font(bold=True)
                        header_cell.fill = header_fill

                    current_row += 1

            for row_num in heading_rows:
                sheet.merge_cells(
                    start_row=row_num, start_column=1,
                    end_row=row_num, end_column=max(max_cols_seen, 1)
                )

            for col_index, width in width_tracker.items():
                col_letter = get_column_letter(col_index)
                sheet.column_dimensions[col_letter].width = min(width + 4, 40)

    if len(workbook.sheetnames) == 0:
        workbook.create_sheet("Sheet1")

    workbook.save(output_path)

    os.remove(input_path)

    return send_file(output_path, as_attachment=True, download_name="converted.xlsx")

def detect_borderless_tables(page, exclude_bboxes, min_rows=5, min_cols=2):
    words = page.get_text("words")
    # Skip pages that are mostly brochure/poster layouts
    images = page.get_images(full=True)

    if len(images) >= 2:
        return []

    def in_excluded(x0, y0, x1, y1):
        cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
        for (bx0, by0, bx1, by1) in exclude_bboxes:
            if bx0 - 2 <= cx <= bx1 + 2 and by0 - 2 <= cy <= by1 + 2:
                return True
        return False

    words = [w for w in words if not in_excluded(w[0], w[1], w[2], w[3])]
    if not words:
        return []

    words_sorted = sorted(words, key=lambda w: (round(w[1], 1), w[0]))
    lines = []
    for w in words_sorted:
        x0, y0, x1, y1, word = w[0], w[1], w[2], w[3], w[4]
        if lines and abs(y0 - lines[-1]["y0"]) < 3:
            lines[-1]["words"].append((x0, y0, x1, y1, word))
            lines[-1]["y1"] = max(lines[-1]["y1"], y1)
        else:
            lines.append({"y0": y0, "y1": y1, "words": [(x0, y0, x1, y1, word)]})

    lines.sort(key=lambda l: l["y0"])

    regions = []
    for line in lines:
        if regions and line["y0"] - regions[-1]["lines"][-1]["y1"] < 15:
            regions[-1]["lines"].append(line)
        else:
            regions.append({"lines": [line]})

    tables = []
    for region in regions:
        rlines = region["lines"]
        if len(rlines) < min_rows:
            continue

        all_x0 = min(w[0] for l in rlines for w in l["words"])
        all_x1 = max(w[2] for l in rlines for w in l["words"])

        bin_size = 1
        n_bins = int((all_x1 - all_x0) / bin_size) + 1
        coverage_count = [0] * n_bins
        for l in rlines:
            for (x0, y0, x1, y1, word) in l["words"]:
                sb = max(0, int((x0 - all_x0) / bin_size))
                eb = min(n_bins, int((x1 - all_x0) / bin_size) + 1)
                for b in range(sb, eb):
                    coverage_count[b] += 1

        threshold = max(1, int(len(rlines) * 0.5))
        covered = [coverage_count[b] >= threshold for b in range(n_bins)]

        gaps = []
        b = 0
        while b < n_bins:
            if not covered[b]:
                start = b
                while b < n_bins and not covered[b]:
                    b += 1
                if (b - start) * bin_size >= 4:
                    gaps.append((all_x0 + start * bin_size, all_x0 + b * bin_size))
            else:
                b += 1

        if len(gaps) < min_cols - 1:
            continue

        col_bounds = [all_x0]
        for g0, g1 in gaps:
            col_bounds.append((g0 + g1) / 2)
        col_bounds.append(all_x1)
        col_bounds = sorted(set(col_bounds))

        if len(col_bounds) - 1 < min_cols:
            continue
        # Reject very narrow columns
        col_widths = [
            col_bounds[i + 1] - col_bounds[i]
            for i in range(len(col_bounds) - 1)
        ]

        if any(w < 25 for w in col_widths):
            continue

        n_test_cols = len(col_bounds) - 1
        
        from collections import Counter

        pattern = []

        for l in rlines:
            cols_hit = set()

            for (x0, y0, x1, y1, word) in l["words"]:
                wcx = (x0 + x1) / 2

                for ci in range(n_test_cols):
                    if col_bounds[ci] <= wcx <= col_bounds[ci + 1]:
                        cols_hit.add(ci)
                        break

            pattern.append(tuple(sorted(cols_hit)))

        most_common = Counter(pattern).most_common(1)[0][1]

        if most_common < len(pattern) * 0.8:
            continue

        cell_count = len(rlines) * (len(col_bounds) - 1)

        if cell_count < 15:
            continue

        # Too much empty horizontal space -> likely brochure/layout
        page_width = page.rect.width
        table_width = all_x1 - all_x0

        if len(col_bounds) <= 3 and len(rlines) <= 6:
            continue

        row_ys = [rlines[0]["y0"]]
        for i in range(len(rlines) - 1):
            row_ys.append((rlines[i]["y1"] + rlines[i + 1]["y0"]) / 2)
        row_ys.append(rlines[-1]["y1"])

        tables.append({
            "bbox": (all_x0, rlines[0]["y0"], all_x1, rlines[-1]["y1"]),
            "row_ys": row_ys,
            "col_xs": col_bounds,
            "v_lines": [],
            "borderless": True,
        })

    return tables    

def detect_tables(drawings):
    h_lines = []
    v_lines = []

    for drawing in drawings:
        for item in drawing["items"]:
            if item[0] != "l":
                continue
            p1, p2 = item[1], item[2]
            x0, y0 = p1.x, p1.y
            x1, y1 = p2.x, p2.y

            if abs(y1 - y0) < 1 and abs(x1 - x0) > 5:
                h_lines.append((round((y0 + y1) / 2, 1), min(x0, x1), max(x0, x1)))
            elif abs(x1 - x0) < 1 and abs(y1 - y0) > 5:
                v_lines.append((round((x0 + x1) / 2, 1), min(y0, y1), max(y0, y1)))

    if len(h_lines) < 2 or len(v_lines) < 2:
        return []

    def cluster(values, tol=2):
        values = sorted(values)
        clusters = []
        for v in values:
            if clusters and v - clusters[-1][-1] <= tol:
                clusters[-1].append(v)
            else:
                clusters.append([v])
        return [sum(c) / len(c) for c in clusters]

    v_sorted = sorted(v_lines, key=lambda v: v[1])
    groups = []
    for v in v_sorted:
        placed = False
        for g in groups:
            gy0, gy1 = g["y_range"]
            if not (v[2] < gy0 - 3 or v[1] > gy1 + 3):
                g["v_lines"].append(v)
                g["y_range"] = (min(gy0, v[1]), max(gy1, v[2]))
                placed = True
                break
        if not placed:
            groups.append({"v_lines": [v], "y_range": (v[1], v[2])})

    tables = []
    for g in groups:
        col_xs = cluster([v[0] for v in g["v_lines"]])
        if len(col_xs) < 2:
            continue

        gy0, gy1 = g["y_range"]
        table_x0 = min(col_xs)
        table_x1 = max(col_xs)

        table_width = table_x1 - table_x0
        relevant_h = [
            h for h in h_lines
            if gy0 - 5 <= h[0] <= gy1 + 5
            and h[2] > table_x0 and h[1] < table_x1
            and (min(h[2], table_x1) - max(h[1], table_x0)) > table_width * 0.15
        ]
        row_ys = cluster([h[0] for h in relevant_h])

        if len(row_ys) < 2:
            continue

        tables.append({
            "bbox": (table_x0, min(row_ys), table_x1, max(row_ys)),
            "row_ys": sorted(row_ys),
            "col_xs": sorted(col_xs),
            "v_lines": g["v_lines"],
            "h_lines": relevant_h,
        })

    return tables
# =========================
# RENDER HELPERS
# =========================

def classify_drawings(drawings):
    """
    Splits page.get_drawings() output into:
      - simple_drawings: items are only 're' (rect) and/or 'l' (line) ->
        stay editable, handled by render_filled_rectangles() / render_lines().
      - complex_drawings: contains at least one 'c' (bezier curve) item ->
        a real curved vector graphic (speech bubble, balloon, icon) that
        python-pptx has no native shape for, so it is rendered as an
        isolated transparent PNG overlay by render_complex_vectors().
    """
    simple_drawings = []
    complex_drawings = []
    for drawing in drawings:
        ops = {item[0] for item in drawing.get("items", [])}
        if "c" in ops:
            complex_drawings.append(drawing)
        else:
            simple_drawings.append(drawing)
    return simple_drawings, complex_drawings


def render_complex_vectors(slide, page, complex_drawings, scale_x, scale_y, zoom=3):
    """
    Renders every complex ('c' / bezier) vector drawing as its own tightly
    cropped, transparent PNG overlay - never the whole page.

    To isolate a single drawing without picking up text/images/other shapes
    that happen to sit in the same bounding box, each drawing's path is
    replayed onto a blank temp page the same size as the real page (using
    PyMuPDF's Shape API with the drawing's original fill/stroke/width), and
    only that temp page is rasterized with alpha=True.

    fill_opacity/stroke_opacity are applied by scaling the rasterized PNG's
    alpha channel directly with PIL, rather than relying on PyMuPDF's own
    opacity support in Shape.finish() (not present in every PyMuPDF
    version, and silently drawing fully opaque otherwise - e.g. a
    decorative balloon's "basket" sub-path can carry a fill color meant to
    be seen only faintly, which would otherwise render as a solid block).
    When a drawing has both a fill and a stroke with different opacities,
    each is rendered as its own pass and composited, so a solid outline
    over a translucent fill (or vice versa) reproduces correctly.

    Call this AFTER simple filled-rectangle backgrounds (render_filled_rectangles)
    so a decorative graphic sitting on a colored panel isn't hidden behind
    it, but BEFORE images/tables/lines/text so it never covers them.
    """

    def _raster_pass(items, page_w, page_h, rect, color, fill, width, even_odd, close_path):
        tmp_doc = fitz.open()
        try:
            tmp_page = tmp_doc.new_page(width=page_w, height=page_h)
            shape = tmp_page.new_shape()
            for item in items:
                op = item[0]
                if op == "c":
                    shape.draw_bezier(item[1], item[2], item[3], item[4])
                elif op == "l":
                    shape.draw_line(item[1], item[2])
                elif op == "re":
                    shape.draw_rect(item[1])
                elif op == "qu":
                    shape.draw_quad(item[1])
            shape.finish(
                color=color, fill=fill, width=width,
                even_odd=even_odd, closePath=close_path,
            )
            shape.commit()
            pix = tmp_page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), clip=rect, alpha=True)
            return Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGBA")
        finally:
            tmp_doc.close()

    def _scale_alpha(img, opacity):
        if opacity >= 1:
            return img
        r, g, b, a = img.split()
        a = a.point(lambda v: round(v * max(opacity, 0)))
        img.putalpha(a)
        return img

    for drawing in complex_drawings:
        rect = drawing.get("rect")
        if rect is None or rect.is_empty or rect.width == 0 or rect.height == 0:
            continue

        items = drawing.get("items", [])
        fill = drawing.get("fill")
        color = drawing.get("color")
        fill_opacity = drawing.get("fill_opacity", 1)
        stroke_opacity = drawing.get("stroke_opacity", 1)
        width = drawing.get("width") or 0
        even_odd = drawing.get("even_odd", False)
        close_path = drawing.get("closePath", True)
        page_w, page_h = page.rect.width, page.rect.height

        if fill is not None and fill_opacity == 0:
            fill = None
        if color is not None and stroke_opacity == 0:
            color = None
        if fill is None and color is None:
            continue  # nothing visible in the original PDF at all

        try:
            if fill is not None and color is not None and fill_opacity != stroke_opacity:
                # Different opacities for fill vs stroke - render each
                # separately so both come out correct, then composite
                # (fill first, stroke on top).
                fill_img = _scale_alpha(
                    _raster_pass(items, page_w, page_h, rect, None, fill, width, even_odd, close_path),
                    fill_opacity,
                )
                stroke_img = _scale_alpha(
                    _raster_pass(items, page_w, page_h, rect, color, None, width, even_odd, close_path),
                    stroke_opacity,
                )
                fill_img.alpha_composite(stroke_img)
                final_img = fill_img
            else:
                final_img = _raster_pass(items, page_w, page_h, rect, color, fill, width, even_odd, close_path)
                opacity = fill_opacity if fill is not None else stroke_opacity
                final_img = _scale_alpha(final_img, opacity)

            buf = io.BytesIO()
            final_img.save(buf, format="PNG")
            png_bytes = buf.getvalue()
        except Exception as e:
            print(f"Complex vector render failed, skipping: {e}")
            png_bytes = None

        if png_bytes is None:
            continue

        left = Emu(round(rect.x0 * scale_x))
        top = Emu(round(rect.y0 * scale_y))
        width = Emu(round(rect.width * scale_x))
        height = Emu(round(rect.height * scale_y))
        if width <= 0 or height <= 0:
            continue

        slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)


def render_filled_rectangles(slide, simple_drawings, table_bboxes, scale_x, scale_y):
    """
    Renders simple single-rectangle filled drawings as editable PowerPoint
    rectangle shapes. Skips any rectangle whose center falls inside a
    detected table bbox, since the table's own cell fill already covers
    that area.

    Respects the drawing's fill_opacity: a rectangle with opacity 0 was
    never actually visible in the PDF (e.g. an invisible hit-box, or a
    "basket" sub-shape meant to be seen only via its outline) and is
    skipped entirely; partial opacity is reproduced with a native OOXML
    alpha value so the shape stays a real, editable, semi-transparent
    rectangle instead of turning solid.
    """
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.oxml.ns import qn

    for drawing in simple_drawings:
        items = drawing.get("items", [])
        if len(items) != 1 or items[0][0] != "re":
            continue

        fill = drawing.get("fill")
        rect = drawing.get("rect")
        if not (fill and rect):
            continue

        fill_opacity = drawing.get("fill_opacity", 1)
        if fill_opacity == 0:
            continue  # never actually visible in the PDF

        cx, cy = (rect.x0 + rect.x1) / 2, (rect.y0 + rect.y1) / 2
        inside_table = any(
            tb[0] - 2 <= cx <= tb[2] + 2 and tb[1] - 2 <= cy <= tb[3] + 2
            for tb in table_bboxes
        )
        if inside_table:
            continue

        left = Emu(round(rect.x0 * scale_x))
        top = Emu(round(rect.y0 * scale_y))
        width = Emu(round((rect.x1 - rect.x0) * scale_x))
        height = Emu(round((rect.y1 - rect.y0) * scale_y))

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
        )
        r, g, b = [int(c * 255) for c in fill]
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()

        if fill_opacity < 1:
            srgb_clr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
            if srgb_clr is not None:
                alpha_el = srgb_clr.makeelement(
                    qn("a:alpha"), {"val": str(int(max(fill_opacity, 0) * 100000))}
                )
                srgb_clr.append(alpha_el)


def find_cell_background_color(simple_drawings, cx0, cy0, cx1, cy1):
    """
    Looks for a simple filled-rectangle drawing whose rect covers most of a
    table cell's area, and returns its RGB fill color. Falls back to None
    (caller should default to white) if no matching rectangle is found.

    This preserves the PDF's own cell highlight colors (e.g. a "Lunch" or
    header row/column drawn with a colored background rectangle behind the
    text) instead of forcing every table cell to plain white.
    """
    cell_area = max((cx1 - cx0) * (cy1 - cy0), 1)
    best_color = None
    best_overlap = 0

    for drawing in simple_drawings:
        items = drawing.get("items", [])
        if len(items) != 1 or items[0][0] != "re":
            continue
        fill = drawing.get("fill")
        rect = drawing.get("rect")
        if not (fill and rect):
            continue

        ox0 = max(cx0, rect.x0)
        oy0 = max(cy0, rect.y0)
        ox1 = min(cx1, rect.x1)
        oy1 = min(cy1, rect.y1)
        if ox1 <= ox0 or oy1 <= oy0:
            continue

        overlap = (ox1 - ox0) * (oy1 - oy0)
        if overlap > best_overlap and overlap >= cell_area * 0.5:
            best_overlap = overlap
            best_color = fill

    return best_color


def render_lines(slide, simple_drawings, table_bboxes, page_rect, scale_x, scale_y):
    """
    Renders 'l' (line) items as editable PowerPoint connector shapes.
    Skips lines inside a detected table bbox (the table grid already draws
    those) and skips near-full-page-width/height lines that are almost
    certainly page-border artifacts rather than intentional content lines.

    Respects stroke_opacity: an opacity-0 line was never actually visible
    in the PDF (skipped entirely) and partial opacity is applied as a
    native OOXML alpha on the connector's stroke color, so a faint
    decorative line (e.g. a lattice pattern meant to blend into a graphic)
    doesn't render as a bold, fully-opaque stroke.
    """
    from pptx.oxml.ns import qn

    for drawing in simple_drawings:
        stroke_opacity = drawing.get("stroke_opacity", 1)
        if stroke_opacity == 0:
            continue

        for item in drawing.get("items", []):
            if item[0] != "l":
                continue

            p1, p2 = item[1], item[2]
            x0, y0 = p1.x, p1.y
            x1, y1 = p2.x, p2.y

            skip = False
            for tb in table_bboxes:
                if tb[0] - 2 <= x0 <= tb[2] + 2 and tb[1] - 2 <= y0 <= tb[3] + 2:
                    skip = True
                    break
            if skip:
                continue

            line_len_x = abs(x1 - x0)
            line_len_y = abs(y1 - y0)

            if line_len_x > page_rect.width * 0.95 and line_len_y < 2:
                continue
            if line_len_y > page_rect.height * 0.95 and line_len_x < 2:
                continue

            left = Emu(round(min(x0, x1) * scale_x))
            top = Emu(round(min(y0, y1) * scale_y))
            width = Emu(max(round(abs(x1 - x0) * scale_x), 1))
            height = Emu(max(round(abs(y1 - y0) * scale_y), 1))

            line_shape = slide.shapes.add_connector(1, left, top, left + width, top + height)

            color = drawing.get("color")
            if color:
                r, g, b = [int(c * 255) for c in color]
                line_shape.line.color.rgb = RGBColor(r, g, b)
            else:
                line_shape.line.color.rgb = RGBColor(0, 0, 0)

            line_width = drawing.get("width") or 1
            line_shape.line.width = Pt(max(line_width, 0.5))

            dashes = drawing.get("dashes")
            if dashes and dashes != "[] 0":
                line_shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

            if stroke_opacity < 1:
                srgb_clr = line_shape.line.color._xFill.find(qn("a:srgbClr"))
                if srgb_clr is not None:
                    alpha_el = srgb_clr.makeelement(
                        qn("a:alpha"), {"val": str(int(max(stroke_opacity, 0) * 100000))}
                    )
                    srgb_clr.append(alpha_el)

def set_cell_borders(cell, color="000000", width_pt=0.75):
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)

        ln = tcPr.makeelement(qn(tag), {})
        ln.set("w", str(int(width_pt * 12700)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {})
        srgbClr.set("val", color)
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = ln.makeelement(qn("a:prstDash"), {})
        prstDash.set("val", "solid")
        ln.append(prstDash)

        tcPr.append(ln)


@app.route("/convert-to-ppt-editable", methods=["POST"])
def convert_to_ppt_editable():
    if "file" not in request.files:
        return {"error": "No file uploaded"}, 400

    file = request.files["file"]

    if file.filename == "":
        return {"error": "Empty filename"}, 400
    temp_dir = tempfile.gettempdir()
    input_path = os.path.join(temp_dir, f"{uuid.uuid4()}.pdf")
    output_path = os.path.join(temp_dir, f"{uuid.uuid4()}.pptx")

    file.save(input_path)

    pdf_doc = fitz.open(input_path)
    prs = Presentation()
    # Use the PDF's own page size for the slide canvas instead of a fixed
    # 16:9 (13.333x7.5in), so a non-widescreen source (e.g. an 11x8.5in
    # tri-fold brochure) doesn't get stretched to fit a different aspect
    # ratio. 1 PDF point = 12700 EMU exactly.
    first_page_rect = pdf_doc[0].rect
    prs.slide_width = Emu(round(first_page_rect.width * 12700))
    prs.slide_height = Emu(round(first_page_rect.height * 12700))
    blank_layout = prs.slide_layouts[6]

    for page in pdf_doc:
        page_rect = page.rect
        # Uniform scale (not independent x/y) so nothing on the page gets
        # stretched non-proportionally, even if a later page's size differs
        # slightly from the first page's.
        scale = min(prs.slide_width / page_rect.width, prs.slide_height / page_rect.height)
        scale_x = scale
        scale_y = scale

        slide = prs.slides.add_slide(blank_layout)

        drawings = page.get_drawings()
        simple_drawings, complex_drawings = classify_drawings(drawings)

        # Table geometry is needed before we draw anything (so rectangles
        # know which regions to skip), but the actual pptx table shapes are
        # built further down - detect_tables()/detect_borderless_tables()
        # only return geometry here, nothing is added to the slide yet.
        tables_info = detect_tables(drawings)
        tables_info += detect_borderless_tables(page, [t["bbox"] for t in tables_info])
        table_bboxes = [t["bbox"] for t in tables_info]

        # Paint order matters: background color panels (simple filled
        # rectangles) go down first, then decorative complex vector
        # graphics (speech bubbles, balloons, icons) on top of them - a
        # balloon drawn behind an opaque background panel would otherwise
        # be completely hidden if the panel were added afterward. Images,
        # tables, and text are added later still, so none of this can
        # cover editable content.
        render_filled_rectangles(slide, simple_drawings, table_bboxes, scale_x, scale_y)
        render_complex_vectors(slide, page, complex_drawings, scale_x, scale_y)

        image_list = page.get_images(full=True)
        processed = set()
        for img in image_list:
            xref = img[0]

            if xref in processed:
                continue
            processed.add(xref)
            try:
                base_image = pdf_doc.extract_image(xref)
                img_bytes = base_image["image"]
                img_ext = base_image["ext"]
                smask_xref = base_image.get("smask", 0)

                if smask_xref:
                    base_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    mask_info = pdf_doc.extract_image(smask_xref)
                    mask_img = Image.open(io.BytesIO(mask_info["image"])).convert("L")
                    if mask_img.size != base_img.size:
                        mask_img = mask_img.resize(base_img.size)
                    base_img.putalpha(mask_img)
                    buf = io.BytesIO()
                    base_img.save(buf, format="PNG")
                    img_bytes = buf.getvalue()
                    img_ext = "png"
            except Exception as e:
                print(f"Image extraction failed: {e}")
                continue

            img_rects = page.get_image_rects(xref)
            if not img_rects:
                continue

            for rect in img_rects:
                # Some PDFs place images that intentionally bleed past the
                # page edge (full-bleed backgrounds). Left as-is, that
                # produces a picture shape whose position/size extends
                # outside the slide canvas - visible as the image
                # "spilling" into the gray pasteboard area in PowerPoint.
                # Clip the placement rect to the page, and crop the actual
                # image bytes to match, so only the part that was ever
                # visible on the page gets embedded.
                vis_x0 = max(rect.x0, page_rect.x0)
                vis_y0 = max(rect.y0, page_rect.y0)
                vis_x1 = min(rect.x1, page_rect.x1)
                vis_y1 = min(rect.y1, page_rect.y1)
                if vis_x1 <= vis_x0 or vis_y1 <= vis_y0:
                    continue  # entirely off-page

                place_bytes = img_bytes
                place_ext = img_ext
                if (vis_x0, vis_y0, vis_x1, vis_y1) != (rect.x0, rect.y0, rect.x1, rect.y1):
                    try:
                        with Image.open(io.BytesIO(img_bytes)) as src_img:
                            src_img.load()
                            px_w, px_h = src_img.size
                            fx = px_w / rect.width
                            fy = px_h / rect.height
                            crop_box = (
                                round((vis_x0 - rect.x0) * fx),
                                round((vis_y0 - rect.y0) * fy),
                                round(px_w - (rect.x1 - vis_x1) * fx),
                                round(px_h - (rect.y1 - vis_y1) * fy),
                            )
                            cropped = src_img.crop(crop_box)
                            buf = io.BytesIO()
                            cropped.save(buf, format="PNG")
                            place_bytes = buf.getvalue()
                            place_ext = "png"
                    except Exception as e:
                        print(f"Off-page image crop failed, using clipped position only: {e}")

                img_path = os.path.join(temp_dir, f"{uuid.uuid4()}.{place_ext}")
                with open(img_path, "wb") as f:
                    f.write(place_bytes)

                left = Emu(round(vis_x0 * scale_x))
                top = Emu(round(vis_y0 * scale_y))
                width = Emu(round((vis_x1 - vis_x0) * scale_x))
                height = Emu(round((vis_y1 - vis_y0) * scale_y))

                try:
                    slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                finally:
                    if os.path.exists(img_path):
                       os.remove(img_path)
        for table_info in tables_info:
            table_bbox = table_info["bbox"]

            # A detected "table" that sits almost entirely inside a
            # decorative complex-vector graphic (e.g. a hot air balloon's
            # basket drawn with a criss-cross lattice of simple lines) is a
            # false positive, not real tabular data - building it as a
            # native table would draw visible cell borders/fill on top of
            # what should be a subtle background graphic. Skip it.
            tb_area = max((table_bbox[2] - table_bbox[0]) * (table_bbox[3] - table_bbox[1]), 1)
            decorative = False
            for cd in complex_drawings:
                cd_rect = cd.get("rect")
                if cd_rect is None:
                    continue
                ox0 = max(table_bbox[0], cd_rect.x0)
                oy0 = max(table_bbox[1], cd_rect.y0)
                ox1 = min(table_bbox[2], cd_rect.x1)
                oy1 = min(table_bbox[3], cd_rect.y1)
                if ox1 <= ox0 or oy1 <= oy0:
                    continue
                overlap = (ox1 - ox0) * (oy1 - oy0)
                if overlap >= tb_area * 0.8:
                    decorative = True
                    break
            if decorative:
                continue

            row_ys = table_info["row_ys"]
            col_xs = table_info["col_xs"]
            v_lines_raw = table_info.get("v_lines", [])

            n_rows = len(row_ys) - 1
            n_cols = len(col_xs) - 1

            if n_rows <= 0 or n_cols <= 0:
                continue

            t_left = Emu(round(col_xs[0] * scale_x))
            t_top = Emu(round(row_ys[0] * scale_y))
            t_width = Emu(round((col_xs[-1] - col_xs[0]) * scale_x))
            t_height = Emu(round((row_ys[-1] - row_ys[0]) * scale_y))

            

            graphic_frame = slide.shapes.add_table(n_rows, n_cols, t_left, t_top, t_width, t_height)
            pptx_table = graphic_frame.table
            pptx_table.first_row = False
            pptx_table.horz_banding = False

            for i in range(n_cols):
                col_width = col_xs[i + 1] - col_xs[i]
                pptx_table.columns[i].width = Emu(round(col_width * scale_x))

            for row in pptx_table.rows:
                for cell in row.cells:
                    cell.margin_left = Pt(0.5)
                    cell.margin_right = Pt(0.5)
                    cell.margin_top = Pt(0.5)
                    cell.margin_bottom = Pt(0.5)

            avg_width = 0.75
            
            for row in pptx_table.rows:
                for cell in row.cells:
                    set_cell_borders(cell, width_pt=avg_width)

            for r_idx, row in enumerate(pptx_table.rows):
                for c_idx, cell in enumerate(row.cells):
                    cell_color = find_cell_background_color(
                        simple_drawings,
                        col_xs[c_idx], row_ys[r_idx],
                        col_xs[c_idx + 1], row_ys[r_idx + 1],
                    )
                    cell.fill.solid()
                    if cell_color:
                        r, g, b = [int(c * 255) for c in cell_color]
                        cell.fill.fore_color.rgb = RGBColor(r, g, b)
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            is_borderless = table_info.get("borderless", False)

            tol = 2
            row_groups = []

            if is_borderless:
                for r in range(n_rows):
                    row_groups.append([(c, c) for c in range(n_cols)])
            else:
              for r in range(n_rows):
                ry0, ry1 = row_ys[r], row_ys[r + 1]
                divider_exists = [False] * (n_cols - 1)
                for i in range(n_cols - 1):
                    bx = col_xs[i + 1]
                    segments = sorted(
                        (vy0, vy1) for (vx, vy0, vy1) in v_lines_raw
                        if abs(vx - bx) <= tol
                    )
                    covered_end = ry0
                    for s0, s1 in segments:
                        if s0 > covered_end + tol:
                            break
                        if s1 > covered_end:
                            covered_end = s1
                    if covered_end >= ry1 - tol:
                        divider_exists[i] = True

                groups = []
                c = 0
                while c < n_cols:
                    start = c
                    while c < n_cols - 1 and not divider_exists[c]:
                        c += 1
                    end = c
                    groups.append((start, end))
                    c += 1
                row_groups.append(groups)

            h_lines_raw = table_info.get("h_lines", [])
            col_row_groups = []
            for c in range(n_cols):
                cx0, cx1 = col_xs[c], col_xs[c + 1]
                v_divider_exists = [False] * (n_rows - 1)
                for j in range(n_rows - 1):
                    by = row_ys[j + 1]
                    segments = sorted(
                        (hx0, hx1) for (hy, hx0, hx1) in h_lines_raw
                        if abs(hy - by) <= tol
                    )
                    covered_end = cx0
                    for s0, s1 in segments:
                        if s0 > covered_end + tol:
                            break
                        if s1 > covered_end:
                            covered_end = s1
                    if covered_end >= cx1 - tol:
                        v_divider_exists[j] = True

                groups = []
                r = 0
                while r < n_rows:
                    start = r
                    while r < n_rows - 1 and not v_divider_exists[r]:
                        r += 1
                    end = r
                    groups.append((start, end))
                    r += 1
                col_row_groups.append(groups)

            for c in range(n_cols):
                for (start, end) in col_row_groups[c]:
                    if end > start:
                        same_h_group = all(
                            row_groups[r] == row_groups[start] for r in range(start, end + 1)
                        )
                        if same_h_group:
                            try:
                                pptx_table.cell(start, c).merge(pptx_table.cell(end, c))
                            except Exception:
                                pass

            for r in range(n_rows):
                for (start, end) in row_groups[r]:
                    if end > start:
                        try:
                            pptx_table.cell(r, start).merge(pptx_table.cell(r, end))
                        except Exception:
                            pass

            def find_col_for_row(row_idx, lx0, lx1):
                best_start, best_overlap = None, 0
                for (start, end) in row_groups[row_idx]:
                    gx0 = col_xs[start]
                    gx1 = col_xs[end + 1]
                    overlap = min(lx1, gx1) - max(lx0, gx0)
                    if overlap > best_overlap:
                        best_overlap = overlap
                        best_start = start
                return best_start

            cell_lines = {}

            page_text_dict = page.get_text("dict")
            for block in page_text_dict["blocks"]:
                if block["type"] != 0:
                    continue
                for line in block["lines"]:
                    line_text = "".join(span["text"] for span in line["spans"]).strip()
                    if not line_text:
                        continue

                    lx0, ly0, lx1, ly1 = line["bbox"]
                    cx = (lx0 + lx1) / 2
                    cy = (ly0 + ly1) / 2

                    if not (table_bbox[0] <= cx <= table_bbox[2] and table_bbox[1] <= cy <= table_bbox[3]):
                        continue

                    best_row, best_overlap = None, 0
                    for i in range(n_rows):
                        ry0, ry1 = row_ys[i], row_ys[i + 1]
                        overlap = min(ly1, ry1) - max(ly0, ry0)
                        if overlap > best_overlap:
                            best_overlap = overlap
                            best_row = i
                    row_idx = best_row

                    if row_idx is None:
                        continue

                    col_idx = find_col_for_row(row_idx, lx0, lx1)
                    if col_idx is None:
                        continue

                    key = (row_idx, col_idx)
                    cell_lines.setdefault(key, []).append((cy, line_text))

            for key in cell_lines:
                entries = sorted(cell_lines[key], key=lambda t: t[0])
                merged = []
                for cy, text in entries:
                    if merged and abs(cy - merged[-1][0]) < 3:
                        merged[-1] = (merged[-1][0], merged[-1][1] + " " + text)
                    else:
                        merged.append((cy, text))
                cell_lines[key] = merged

            align_mode = PP_ALIGN.LEFT if n_cols <= 6 else PP_ALIGN.CENTER

            for (row_idx, col_idx), lines in cell_lines.items():
                lines.sort(key=lambda t: t[0])
                cell = pptx_table.cell(row_idx, col_idx)
                full_text = "\n".join(t[1] for t in lines)
                cell.text_frame.text = full_text
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                col_width_pt = (col_xs[col_idx + 1] - col_xs[col_idx])
                longest_line = max((len(t[1]) for t in lines), default=1)
                est_char_width = col_width_pt / max(longest_line, 1)
                adaptive_size = max(min(int(est_char_width * 1.8), 9), 5)

                for para in cell.text_frame.paragraphs:
                    para.alignment = align_mode
                    for run in para.runs:
                        run.font.size = Pt(adaptive_size)
                        run.font.color.rgb = RGBColor(0, 0, 0)
        # Lines only here (simple filled rectangles already went down
        # earlier, before images/tables, so background panels don't cover
        # decorative overlays). Lines still wait until after tables so they
        # can skip anything inside a detected table's bbox.
        render_lines(slide, simple_drawings, table_bboxes, page_rect, scale_x, scale_y)

        page_text_dict = page.get_text("dict")
        text_blocks = [b for b in page_text_dict["blocks"] if b["type"] == 0]

        all_lines = []
        for block in text_blocks:
            for line in block["lines"]:
                has_text = any(span["text"].strip() for span in line["spans"])
                if not has_text:
                    continue

                lx0, ly0, lx1, ly1 = line["bbox"]
                lcx = (lx0 + lx1) / 2
                lcy = (ly0 + ly1) / 2

                skip = False
                for tb in table_bboxes:
                    if (tb[0] - 2 <= lcx <= tb[2] + 2 and
                        tb[1] - 2 <= lcy <= tb[3] + 2):
                        skip = True
                        break
                if skip:
                    continue

                all_lines.append({
                    "spans": line["spans"],
                    "bbox": line["bbox"],
                })

        for line in all_lines:
            x0, y0, x1, y1 = line["bbox"]
            spans_list = [line["spans"]]

            line_height = y1 - y0
            bbox_font_size = line_height * 0.75

            max_final_size = 6
            for spans in spans_list:
                for span in spans:
                    if not span.get("text"):
                        continue
                    font_size = span.get("size", 18)
                    final_size = max(font_size, bbox_font_size)
                    max_final_size = max(max_final_size, final_size)

            left = Emu(round(x0 * scale_x))
            top = Emu(round(y0 * scale_y))
            width = Emu(round((x1 - x0) * scale_x))

            raw_height_pt = (y1 - y0)
            min_height_pt = max_final_size * 1.25
            box_height_pt = max(raw_height_pt, min_height_pt)
            height = Emu(round(box_height_pt * scale_y))

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            # python-pptx textboxes default to ~0.05in top/bottom and
            # ~0.1in left/right internal margins. Left un-zeroed, every
            # line of extracted text renders a few points lower/right than
            # its real PDF position - enough for a closely-spaced divider
            # line (e.g. a rule right under a heading) to visually cut
            # through the text below it. Zeroing this aligns text exactly
            # with the coordinates already computed from the PDF.
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0

            first_para = True
            for spans in spans_list:
                if first_para:
                    para = tf.paragraphs[0]
                    first_para = False
                else:
                    para = tf.add_paragraph()

                for span in spans:
                    text = span["text"]
                    if not text:
                        continue

                    has_pua = any(0xE000 <= ord(ch) <= 0xF8FF for ch in text)
                    if has_pua:
                        cleaned = ""
                        for ch in text:
                            if 0xE000 <= ord(ch) <= 0xF8FF:
                                cleaned += "\u2022 "
                            else:
                                cleaned += ch
                        text = cleaned

                    run = para.add_run()
                    run.text = text

                    font_size = span.get("size", 18)
                    final_size = max(font_size, bbox_font_size)
                    run.font.size = Pt(max(final_size, 6))

                    font_name = span.get("font")
                    if font_name and not has_pua:
                        run.font.name = font_name

                    flags = span.get("flags", 0)
                    run.font.bold = bool(flags & 2 ** 4)
                    run.font.italic = bool(flags & 2 ** 1)

                    color_int = span.get("color", 0)
                    r = (color_int >> 16) & 255
                    g = (color_int >> 8) & 255
                    b = color_int & 255
                    run.font.color.rgb = RGBColor(r, g, b)

    pdf_doc.close()
    prs.save(output_path)
    os.remove(input_path)

    return send_file(output_path, as_attachment=True, download_name="converted.pptx")



if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)